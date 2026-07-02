import os
import json
from pathlib import Path
import fitz  # PyMuPDF
from pypdf import PdfReader # Запасной парсер для битых PDF
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import docx2txt

RAW_DIR = Path("data/raw")
PROCESSED_DIR = Path("data/processed")
PROCESSED_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_FILE = PROCESSED_DIR / "chunks.jsonl"

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1500,
    chunk_overlap=300,
    separators=["\n\n", "\n", ".", " ", ""]
)

def extract_from_pdf(file_path):
    text = ""
    # Попытка 1: Быстрый PyMuPDF
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text
    except Exception as e:
        print(f"⚠️ Ошибка PyMuPDF ({file_path.name}): {e}. Пробую через pypdf...")
    
    # Попытка 2: Запасной pypdf (если PyMuPDF рухнул из-за слоев/шрифтов)
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    except Exception as e:
        print(f"❌ Полный провал с PDF ({file_path.name}): {e}")
        
    return text

def extract_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"❌ Ошибка PPTX ({file_path.name}): {e}")
    return text

def extract_from_word(file_path):
    text = ""
    # Попытка 1: python-docx (идеально сохраняет структуру таблиц)
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    text += " | ".join(row_data) + "\n"
        return text
    except Exception as e:
        print(f"⚠️ Ошибка python-docx ({file_path.name}): {e}. Пробую через docx2txt...")
    
    # Попытка 2: docx2txt (глухой парсер, игнорирует макросы и читает сырой XML)
    try:
        text = docx2txt.process(file_path)
    except Exception as e:
        print(f"❌ Полный провал с Word ({file_path.name}): {e}")
        
    return text if text else ""
def process_documents():
    # Ищем все поддерживаемые форматы
    extensions = ('*.pdf', '*.pptx', '*.docx', '*.docm')
    all_files = []
    for ext in extensions:
        all_files.extend(RAW_DIR.rglob(ext))
        all_files.extend(RAW_DIR.rglob(ext.upper())) # Учитываем капс (напр. .PDF)
        
    # Убираем дубликаты на всякий случай
    all_files = list(set(all_files))
    
    print(f"Найдено файлов для обработки: {len(all_files)}")
    total_chunks = 0
    
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f_out:
        for file_path in all_files:
            ext = file_path.suffix.lower()
            raw_text = ""
            
            if ext == '.pdf':
                raw_text = extract_from_pdf(file_path)
            elif ext == '.pptx':
                raw_text = extract_from_pptx(file_path)
            elif ext in ['.docx', '.docm']:
                raw_text = extract_from_word(file_path)
                
            if not raw_text.strip():
                continue
                
            chunks = text_splitter.split_text(raw_text)
            for chunk in chunks:
                record = {
                    "source": str(file_path.name),
                    "text": chunk
                }
                f_out.write(json.dumps(record, ensure_ascii=False) + "\n")
                total_chunks += 1
                
            print(f"✅ Обработан: {file_path.name} -> {len(chunks)} чанков")
            
    print(f"\nГОТОВО! Всего создано чанков: {total_chunks}")
    print(f"Результат сохранен в: {OUTPUT_FILE}")

if __name__ == "__main__":
    process_documents()